import os
import json
import uuid
import platform
import subprocess
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import openai
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure OpenAI client
openai.api_key = os.getenv("OPENAI_API_KEY")

class PresentationGenerator:
    def __init__(self):
        self.output_dir = "generated"
        
        # Create output directory if it doesn't exist
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)
            
        # Set default parameters
        self.default_model = "gpt-4"
        self.use_llm_chaining = False
    
    def generate_presentation_outline(self, prompt):
        """Use OpenAI to generate a structured presentation outline"""
        try:
            # Using the new client.chat.completions.create format for OpenAI >= 1.0.0
            response = openai.ChatCompletion.create(
                model=self.default_model,
                messages=[
                    {"role": "system", "content": """You are an expert presentation designer. 
                    Create a detailed presentation outline with exactly the following JSON structure:
                    {
                        "title": "Main Presentation Title",
                        "slides": [
                            {
                                "title": "Slide Title",
                                "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
                                "notes": "Optional presenter notes"
                            },
                            ...more slides...
                        ]
                    }
                    
                    Follow these guidelines:
                    - Create 5-10 slides depending on the topic complexity
                    - Include an introduction, main content slides, and conclusion
                    - Keep bullet points concise and focused (1-2 lines each)
                    - Add presenter notes to provide additional context
                    - Ensure logical flow between slides
                    - ONLY output valid JSON that matches the structure above
                    """},
                    {"role": "user", "content": f"Generate a presentation outline about: {prompt}"}
                ],
                temperature=0.7
            )
            
            content = response.choices[0].message.content
            
            # Extract JSON from the response (in case there's surrounding text)
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            json_str = content[json_start:json_end]
            
            # Parse JSON
            outline = json.loads(json_str)
            return outline
            
        except Exception as e:
            error_message = str(e)
            print(f"Error generating outline: {error_message}")
            
            # Check for API quota exceeded error - propagate it instead of using fallback
            if 'insufficient_quota' in error_message or 'exceeded your current quota' in error_message:
                raise Exception(f"OpenAI API quota exceeded: {error_message}")
                
            # Use fallback outline only for non-quota errors
            return {
                "title": f"Presentation about {prompt}",
                "slides": [
                    {
                        "title": "Introduction", 
                        "content": ["Overview of the topic", "Key points to be covered"],
                        "notes": "Introduce yourself and the topic"
                    },
                    {
                        "title": "Key Points",
                        "content": ["Main idea 1", "Main idea 2", "Main idea 3"],
                        "notes": "Explain the main concepts"
                    },
                    {
                        "title": "Conclusion",
                        "content": ["Summary of key points", "Call to action or next steps"],
                        "notes": "Wrap up and invite questions"
                    }
                ]
            }
    
    def research_slide_content(self, slide_info, main_topic):
        """Research and enhance content for a specific slide"""
        try:
            # Prepare the prompt for research
            slide_title = slide_info["title"]
            current_content = " ".join(slide_info["content"])
            
            research_prompt = f"""
            I'm creating a presentation slide titled "{slide_title}" as part of a presentation about "{main_topic}".
            
            The current content for this slide is:
            {current_content}
            
            Please research this topic in depth and provide:
            1. 4-6 bullet points with specific facts, data, or examples
            2. Detailed presenter notes with background information and talking points
            3. Any relevant statistics or case studies that would strengthen this slide
            
            Format your response as a JSON object with these fields:
            {{
                "bullet_points": ["Point 1 with specific data", "Point 2 with example", ...],
                "presenter_notes": "Detailed background information and talking points for the presenter",
                "references": ["Optional reference 1", "Optional reference 2", ...]
            }}
            """
            
            response = openai.ChatCompletion.create(
                model=self.default_model,
                messages=[
                    {"role": "system", "content": "You are a research assistant specializing in creating well-researched, factual presentation content."},
                    {"role": "user", "content": research_prompt}
                ],
                temperature=0.5
            )
            
            content = response.choices[0].message.content
            
            # Extract JSON from the response
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            json_str = content[json_start:json_end]
            
            # Parse JSON
            enhanced_content = json.loads(json_str)
            
            # Update the slide information with the researched content (will be condensed later)
            updated_slide = slide_info.copy()
            updated_slide["researched_content"] = enhanced_content
            updated_slide["content"] = enhanced_content["bullet_points"]
            updated_slide["notes"] = enhanced_content["presenter_notes"]
            
            # Add references if available
            if "references" in enhanced_content and enhanced_content["references"]:
                updated_slide["references"] = enhanced_content["references"]
            
            return updated_slide
            
        except Exception as e:
            error_message = str(e)
            print(f"Error researching slide content: {error_message}")
            
            # Check for API quota exceeded error - propagate it
            if 'insufficient_quota' in error_message or 'exceeded your current quota' in error_message:
                raise Exception(f"OpenAI API quota exceeded: {error_message}")
                
            # Return the original slide for other errors
            return slide_info
    
    def condense_slide_content(self, researched_slide):
        """Condense the researched content to be more concise and focused on essential information"""
        try:
            # Check if this slide has researched content
            if "researched_content" not in researched_slide:
                return researched_slide
                
            slide_title = researched_slide["title"]
            
            # Get the researched content
            bullet_points = researched_slide["content"]
            presenter_notes = researched_slide["notes"]
            references = researched_slide.get("references", [])
            
            # Format the bullet points and notes for the prompt
            bullet_points_text = "\n".join([f"- {point}" for point in bullet_points])
            references_text = "\n".join([f"- {ref}" for ref in references]) if references else "None"
            
            condense_prompt = f"""
            I need to condense the following researched content for a slide titled "{slide_title}".
            
            ORIGINAL BULLET POINTS:
            {bullet_points_text}
            
            ORIGINAL PRESENTER NOTES:
            {presenter_notes}
            
            REFERENCES:
            {references_text}
            
            Please create a more concise version that captures only the most essential information:
            1. Maximum 3-4 bullet points of 10-15 words each
            2. Each bullet point should be a clear, direct statement with specific information
            3. Presenter notes should be 2-3 sentences maximum
            4. Include only the single most important reference, if any
            
            Format your response as a JSON object with these fields:
            {{
                "concise_bullet_points": ["Concise point 1", "Concise point 2", ...],
                "concise_notes": "Very brief presenter notes (2-3 sentences max)",
                "key_reference": "Most important reference (if applicable)"
            }}
            """
            
            response = openai.ChatCompletion.create(
                model=self.default_model,
                messages=[
                    {"role": "system", "content": "You are a specialist in creating extremely concise, information-dense presentation content. Your goal is maximum information in minimum words."},
                    {"role": "user", "content": condense_prompt}
                ],
                temperature=0.3
            )
            
            content = response.choices[0].message.content
            
            # Extract JSON from the response
            json_start = content.find('{')
            json_end = content.rfind('}') + 1
            json_str = content[json_start:json_end]
            
            # Parse JSON
            condensed_content = json.loads(json_str)
            
            # Update the slide with condensed content
            condensed_slide = researched_slide.copy()
            condensed_slide["content"] = condensed_content["concise_bullet_points"]
            condensed_slide["notes"] = condensed_content["concise_notes"]
            
            # Add the key reference if available
            if "key_reference" in condensed_content and condensed_content["key_reference"] and condensed_content["key_reference"].lower() != "none":
                condensed_slide["notes"] += f"\n\nKey reference: {condensed_content['key_reference']}"
            
            return condensed_slide
            
        except Exception as e:
            print(f"Error condensing slide content: {e}")
            # Return the original researched slide if condensing fails
            return researched_slide
    
    def enhance_presentation_outline(self, outline, main_topic):
        """Enhance the presentation outline with researched and condensed content for each slide"""
        enhanced_outline = {
            "title": outline["title"],
            "slides": []
        }
        
        print("Researching and enhancing slide content...")
        total_slides = len(outline["slides"])
        
        for i, slide in enumerate(outline["slides"]):
            print(f"Researching slide {i+1}/{total_slides}: {slide['title']}...")
            # Step 1: Research the slide content
            researched_slide = self.research_slide_content(slide, main_topic)
            
            # Step 2: Condense the researched content
            print(f"Condensing slide {i+1}/{total_slides}: {slide['title']}...")
            condensed_slide = self.condense_slide_content(researched_slide)
            
            enhanced_outline["slides"].append(condensed_slide)
            
            # Add a small delay to avoid rate limiting
            if i < total_slides - 1:
                time.sleep(0.5)
        
        return enhanced_outline
    
    def create_presentation(self, outline):
        """Create a PowerPoint presentation from the structured outline"""
        # Create a new presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = outline["title"]
        subtitle.text = "Presentation Generation"
        
        # Apply some basic styling to title slide
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        
        # Add footer to title slide
        self._add_footer(slide, "Made with ❤️ by Arman")
        
        # Add content slides
        for slide_info in outline["slides"]:
            content_slide_layout = prs.slide_layouts[1]  # Layout with title and content
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set slide title
            title = slide.shapes.title
            title.text = slide_info["title"]
            
            # Add content as bullet points
            content = slide.placeholders[1]
            text_frame = content.text_frame
            
            for i, bullet_point in enumerate(slide_info["content"]):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = bullet_point
                p.level = 0
            
            # Add notes if provided
            if "notes" in slide_info and slide_info["notes"]:
                slide.notes_slide.notes_text_frame.text = slide_info["notes"]
                
            # Add footer to content slide
            self._add_footer(slide, "Made with ❤️ by Arman")
        
        # Generate a unique filename
        file_id = str(uuid.uuid4())[:8]
        file_name = f"{self.output_dir}/presentation_{file_id}.pptx"
        
        # Save the presentation
        prs.save(file_name)
        return file_name
    
    def _add_footer(self, slide, footer_text):
        """Add a footer to a slide"""
        # Add a textbox at the bottom of the slide for the footer
        left = Inches(0.1)
        top = Inches(7.0)
        width = Inches(9.8)
        height = Inches(0.3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        p = tf.paragraphs[0]
        p.text = footer_text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(100, 100, 100)  # Gray color
    
    def convert_to_pdf(self, pptx_path):
        """Convert PowerPoint to PDF using appropriate method based on OS"""
        pdf_path = pptx_path.replace('.pptx', '.pdf')
        
        try:
            system = platform.system()
            
            if system == 'Darwin':  # macOS
                # Check if LibreOffice is installed
                try:
                    # Attempt to convert using LibreOffice
                    subprocess.run([
                        '/Applications/LibreOffice.app/Contents/MacOS/soffice',
                        '--headless',
                        '--convert-to', 'pdf',
                        '--outdir', os.path.dirname(pptx_path),
                        pptx_path
                    ], check=True)
                    print(f"PDF created with LibreOffice: {pdf_path}")
                    return pdf_path
                except (subprocess.SubprocessError, FileNotFoundError):
                    print("LibreOffice not found or error occurred.")
                    
                # Try using unoconv if LibreOffice direct method failed
                try:
                    subprocess.run(['unoconv', '-f', 'pdf', pptx_path], check=True)
                    print(f"PDF created with unoconv: {pdf_path}")
                    return pdf_path
                except (subprocess.SubprocessError, FileNotFoundError):
                    print("unoconv not found or error occurred.")
                    
            elif system == 'Windows':
                # On Windows, use comtypes to convert
                try:
                    import comtypes.client
                    
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = True
                    
                    deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path))
                    deck.SaveAs(os.path.abspath(pdf_path), 32)  # 32 is the PDF format code
                    deck.Close()
                    powerpoint.Quit()
                    
                    print(f"PDF created with PowerPoint: {pdf_path}")
                    return pdf_path
                except Exception as e:
                    print(f"Error using comtypes: {e}")
            
            elif system == 'Linux':
                # On Linux, use LibreOffice
                try:
                    subprocess.run([
                        'libreoffice', 
                        '--headless', 
                        '--convert-to', 'pdf', 
                        '--outdir', os.path.dirname(pptx_path),
                        pptx_path
                    ], check=True)
                    print(f"PDF created with LibreOffice: {pdf_path}")
                    return pdf_path
                except (subprocess.SubprocessError, FileNotFoundError):
                    print("LibreOffice not found or error occurred.")
                    
                # Try using unoconv as fallback
                try:
                    subprocess.run(['unoconv', '-f', 'pdf', pptx_path], check=True)
                    print(f"PDF created with unoconv: {pdf_path}")
                    return pdf_path
                except (subprocess.SubprocessError, FileNotFoundError):
                    print("unoconv not found or error occurred.")
            
            print(f"Could not convert to PDF on {system} system.")
            return None
            
        except Exception as e:
            print(f"Error during PDF conversion: {e}")
            return None
    
    def generate(self, prompt, convert_to_pdf=True, use_llm_chaining=False):
        """Generate a complete presentation from a prompt"""
        self.use_llm_chaining = use_llm_chaining
        
        # Step 1: Generate initial outline
        print("Step 1: Generating presentation outline...")
        outline = self.generate_presentation_outline(prompt)
        
        # Step 2 (Optional): If LLM chaining is enabled, research and enhance each slide
        if use_llm_chaining:
            print("Step 2: Researching and enhancing slide content with concise information...")
            outline = self.enhance_presentation_outline(outline, prompt)
        
        # Step 3: Create the PowerPoint presentation
        print("Step 3: Creating PowerPoint presentation...")
        pptx_path = self.create_presentation(outline)
        
        # Step 4: Convert to PDF if requested
        pdf_path = None
        if convert_to_pdf:
            print("Step 4: Converting to PDF...")
            pdf_path = self.convert_to_pdf(pptx_path)
        
        # Return both paths in a dictionary
        return {
            'pptx_path': pptx_path,
            'pdf_path': pdf_path
        }


if __name__ == "__main__":
    # Test the generator
    generator = PresentationGenerator()
    prompt = "The future of artificial intelligence in healthcare"
    result = generator.generate(prompt, use_llm_chaining=True)
    
    print(f"PowerPoint created: {result['pptx_path']}")
    if result['pdf_path']:
        print(f"PDF created: {result['pdf_path']}")
    else:
        print("PDF conversion not successful. You may need to install LibreOffice or PowerPoint.") 